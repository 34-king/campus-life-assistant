from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLUE = RGBColor(0x15, 0x65, 0xC0)
DARK = RGBColor(0x1A, 0x23, 0x7E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
GRAY = RGBColor(0x66, 0x66, 0x66)

def add_slide(title, content_lines, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    # Title bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    tf = bar.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Content
    left = Inches(0.8)
    top = Inches(1.3)
    width = Inches(11.7)
    height = Inches(5.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16) if not line.startswith("##") else Pt(20)
        p.font.bold = line.startswith("##")
        p.font.color.rgb = DARK if line.startswith("##") else GRAY
        p.space_after = Pt(6)
        if line.startswith("##"):
            p.space_before = Pt(12)
    
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide

# ===== Slide 1: Title =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
# Background
bg = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = BLUE
# Title
txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "校园生活助手"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "移动应用开发实践课程 期末项目"
p.font.size = Pt(24)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER
p.space_before = Pt(12)
# Info
txBox2 = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11), Inches(1.5))
tf2 = txBox2.text_frame
for txt in ["河南理工大学 · 2026年春季", "指导老师：任建吉", "Flutter 3.x + Dart 3.11 + Material Design 3"]:
    p = tf2.add_paragraph()
    p.text = txt
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

# ===== Slide 2: 项目概述 =====
add_slide("项目概述", [
    "## 项目简介",
    "一款面向高校学生的跨平台校园生活助手App",
    "整合校园日常功能：课程、食堂、公告、校园卡、成绩等",
    "",
    "## 技术亮点",
    "• Flutter 3.41 跨平台框架，一套代码支持Android/iOS",
    "• Material Design 3 设计规范，原生级UI体验",
    "• Provider 状态管理，组件化架构",
    "• Trae AI IDE 全程辅助开发（Vibe Coding模式）",
    "",
    "## 核心数据",
    "• 12个功能页面  • 6个数据模型  • 2个Provider  • 140+MB APK"
])

# ===== Slide 3: 功能展示 =====
add_slide("功能展示（一）", [
    "## 首页",
    "• 实时天气（wttr.in API 网络数据交互）",
    "• 今日课程自动展示",
    "• 6个快捷入口（食堂/公告/地图/校园卡/成绩/校历）",
    "",
    "## 课程表",
    "• 周视图网格布局，横向+纵向可滚动",
    "• 按天+节次排列，彩色课程卡片",
    "• 支持跨节次课程显示"
])

# ===== Slide 4: 功能展示2 =====
add_slide("功能展示（二）", [
    "## 食堂菜单",
    "• TabBar 按食堂分类筛选",
    "• 菜品卡片：名称/价格/评分/食堂",
    "• 收藏功能（Provider + SharedPreferences 持久化）",
    "",
    "## 校园公告",
    "• 列表展示公告，类别标签+日期",
    "• 底部弹窗查看详情（DraggableScrollableSheet）",
    "• 公告收藏管理",
    "",
    "## 校园卡",
    "• 虚拟校园卡（余额+学生信息）",
    "• 交易流水列表",
    "• 快捷操作按钮"
])

# ===== Slide 5: 功能展示3 =====
add_slide("功能展示（三）", [
    "## 成绩查询",
    "• 按学期查看成绩",
    "• GPA统计（加权平均计算）",
    "• 分数环形指示器+等级标签",
    "",
    "## 校历",
    "• 月视图日历，月份导航",
    "• 假期/考试/活动事件标记",
    "• 点击日期查看当日事件",
    "",
    "## 校园地图",
    "• 8个主要建筑网格展示",
    "• 点击查看建筑详情",
    "",
    "## 个人中心 / 设置 / 收藏管理"
])

# ===== Slide 6: 技术架构 =====
add_slide("技术架构", [
    "## 整体架构",
    "UI层：12个Page Widget（Material Design 3）",
    "状态管理层：Provider（ChangeNotifierProvider）",
    "   ├─ FavoritesProvider：收藏状态管理",
    "   └─ ThemeProvider：主题模式/通知/字体",
    "数据层：MockData + SharedPreferences + HTTP API",
    "",
    "## 数据流",
    "用户操作 → Provider.notifyListeners() → Consumer Widget → UI更新",
    "",
    "## 路由设计",
    "/splash → /home → 底部导航5Tab + 6个路由页面",
    "（/map, /card, /grade, /calendar, /favorites, /settings）"
])

# ===== Slide 7: 深色模式 =====
add_slide("深色模式与本地持久化", [
    "## 深色模式",
    "• 支持三种模式：跟随系统 / 浅色 / 深色",
    "• Material Design 3 自动适配暗色主题",
    "• 通过 ThemeProvider 管理状态",
    "• 设置页面一键切换",
    "",
    "## 本地数据持久化",
    "• SharedPreferences 存储：",
    "  - 收藏的菜品和公告（List<String>）",
    "  - 主题模式选择（system/light/dark）",
    "  - 消息通知开关",
    "  - 字体大小设置",
    "• App 重启后自动恢复用户设置"
])

# ===== Slide 8: AI协作 =====
add_slide("AI 协作开发", [
    "## 开发工具：Trae AI IDE",
    "",
    "## 协作流程",
    "1. 需求描述 → 用自然语言描述功能",
    "2. AI 生成 → Trae AI Chat 生成初始代码",
    "3. 人工审查 → 检查逻辑，调整UI",
    "4. 整合测试 → 运行验证",
    "5. 迭代优化 → 重复上述步骤",
    "",
    "## 效率提升",
    "• 页面框架代码：AI生成耗时约30秒/页",
    "• 数据处理逻辑：人工优化为主",
    "• UI细节调整：手动微调材质参数",
    "• 总体效率：较传统开发提升约3-5倍"
])

# ===== Slide 9: 课程要求对照 =====
add_slide("课程要求对照", [
    "✅ Flutter 3.x 跨平台开发  →  3.41.9，支持Android/iOS/Web",
    "✅ Android设备正常运行  →  已生成可安装APK（API 21+）",
    "✅ 3-5个核心页面  →  实际12个页面",
    "✅ Material Design 3  →  useMaterial3: true",
    "✅ 深色模式切换  →  三种主题模式",
    "✅ 本地数据持久化  →  SharedPreferences",
    "✅ 网络数据交互  →  wttr.in天气API（HTTP+JSON）",
    "✅ Provider状态管理  →  2个Provider",
    "✅ AI辅助开发  →  Trae AI IDE全程辅助",
    "✅ Git版本控制  →  本仓库管理",
    "✅ 可安装APK  →  桌面校园生活助手.apk",
    "✅ 鸿蒙适配方案  →  见技术报告文档"
])

# ===== Slide 10: 总结 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = GREEN
txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(3))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "谢谢！"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "欢迎各位老师批评指正"
p.font.size = Pt(24)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER
p.space_before = Pt(20)

# Save
output_path = os.path.expanduser("~/Desktop/校园生活助手-答辩PPT.pptx")
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
